"""
Create a 3-slide PPTX template and an Excel workbook with 3 sheets (Appreciation, Recognition, Participation)
for testing the multi-sheet / multi-template feature.

Run: python create_multi_template.py
Produces: multi_template.pptx, multi_data.xlsx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
import pandas as pd

SLIDE_INFO = [
    ("Appreciation", 'Certificate of Appreciation'),
    ("Recognition", 'Certificate of Recognition'),
    ("Participation", 'Certificate of Participation'),
]


def create_pptx(path: str):
    prs = Presentation()
    # remove default slide if present
    while prs.slides:
        try:
            prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])
        except Exception:
            break

    for sheet_name, title_text in SLIDE_INFO:
        slide_layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(slide_layout)

        left = Inches(1)
        top = Inches(1)
        width = Inches(8)
        height = Inches(1.5)

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(36)
        p.font.bold = True

        p = tf.add_paragraph()
        p.text = f'<<FirstName>> <<Surname>>'
        p.font.size = Pt(28)

        # include the sheet name somewhere so the auto-detect can find it
        left2 = Inches(1)
        top2 = Inches(4.5)
        tb = slide.shapes.add_textbox(left2, top2, Inches(8), Inches(0.5))
        tb.text_frame.text = sheet_name

    prs.save(path)
    print('Wrote sample multi-slide template to', path)


def create_excel(path: str):
    # create a dict of DataFrames for three sheets
    data = {
        'Appreciation': pd.DataFrame([
            {'FirstName': 'Alice', 'Surname': 'Anderson'},
            {'FirstName': 'Bob', 'Surname': 'Brown'},
        ]),
        'Recognition': pd.DataFrame([
            {'FirstName': 'Carol', 'Surname': 'Clark'},
        ]),
        'Participation': pd.DataFrame([
            {'FirstName': 'Dave', 'Surname': 'Dawson'},
            {'FirstName': 'Eve', 'Surname': 'Evans'},
            {'FirstName': 'Frank', 'Surname': 'Frost'},
        ]),
    }
    with pd.ExcelWriter(path) as writer:
        for sheet_name, df in data.items():
            df.to_excel(writer, sheet_name=sheet_name, index=False)
    print('Wrote sample multi-sheet data to', path)


if __name__ == '__main__':
    create_pptx('multi_template.pptx')
    create_excel('multi_data.xlsx')

