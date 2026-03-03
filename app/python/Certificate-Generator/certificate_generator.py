"""
Certificate Generator CLI

Usage: python certificate_generator.py --template template.pptx --data data.xlsx --outdir out --format png

Reads a PPTX template with placeholders in the form <<ColumnName>> and an Excel file where each row is a recipient.
Generates a file per row in the requested output format (pptx, pdf, png, jpg). Filename may be customized with a template string
using column names in braces, e.g. "{Initials}_{Surname}_{date}.png". Default is Fullname_date_timestamp.

Notes:
- Requires python-pptx, pandas, pillow, pdf2image.
- For PDF/PNG/JPG conversion, LibreOffice (soffice) and Poppler may be required. See README.
"""
import argparse
import json
import logging
import os
import pathlib
import re
import shutil
import subprocess
import sys
from datetime import datetime
from tempfile import mkstemp, mkdtemp
from typing import List, Dict, Optional, Tuple
from copy import deepcopy

import pandas as pd
from pptx import Presentation
from PIL import Image
from pdf2image import convert_from_path

PLACEHOLDER_RE = re.compile(r"<<([^<>]+)>>")


def parse_mapping_arg(map_arg: Optional[str]) -> Dict[str, int]:
    """Parse mapping string like "SheetA:0,SheetB:2" into dict{"SheetA":0,...}.
    Accepts comma or semicolon separators. Slide indices are 0-based.
    """
    mapping: Dict[str, int] = {}
    if not map_arg:
        return mapping
    parts = [p.strip() for p in re.split(r'[;,]', map_arg) if p.strip()]
    for part in parts:
        if ':' not in part:
            continue
        key, val = part.split(':', 1)
        key = key.strip()
        try:
            idx = int(val.strip())
            mapping[key] = idx
        except ValueError:
            # ignore non-integer mapping
            continue
    return mapping


def find_slide_indices_by_sheetname(prs: Presentation, sheet_name: str) -> List[int]:
    """Return list of slide indices where any text contains the sheet_name (case-insensitive)."""
    found: List[int] = []
    needle = str(sheet_name).strip().lower()
    if not needle:
        return found
    for idx, slide in enumerate(prs.slides):
        matched = False
        for shape in slide.shapes:
            try:
                if not shape.has_text_frame:
                    continue
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        if run.text and needle in run.text.lower():
                            matched = True
                            break
                    if matched:
                        break
            except Exception:
                continue
            if matched:
                break
        if matched:
            found.append(idx)
    return found


def remove_slide_at_index(prs: Presentation, index: int) -> None:
    """Remove the slide at `index` from `prs` by manipulating the slideIdList."""
    # Remove the slide id from the presentation sldIdLst
    sldIdLst = prs.slides._sldIdLst
    if index < 0 or index >= len(sldIdLst):
        return
    sldId = sldIdLst[index]
    sldIdLst.remove(sldId)


def sanitize_filename(name: str) -> str:
    # remove characters invalid on Windows filenames
    invalid = '<>:"/\\|?*'
    for ch in invalid:
        name = name.replace(ch, "_")
    name = name.strip()
    return name


# NEW: ensure we don't overwrite when multiple rows resolve to the same base name
def ensure_unique_basename(out_dir: str, base: str, ext: str) -> str:
    """Return a basename that is unique within out_dir for the given ext, appending (2), (3), ... if needed."""
    candidate = sanitize_filename(str(base))
    path = os.path.join(out_dir, f"{candidate}.{ext}")
    if not os.path.exists(path):
        return candidate
    i = 2
    while True:
        alt = f"{candidate} ({i})"
        path = os.path.join(out_dir, f"{alt}.{ext}")
        if not os.path.exists(path):
            return alt
        i += 1


def resolve_placeholder_value(mapping: Dict[str, object], key: str) -> Tuple[bool, object]:
    normalized_key = str(key or '').strip().lower()
    if normalized_key in mapping:
        return True, mapping.get(normalized_key, '')

    alias_map = {
        'name': ['fullname', 'full name'],
        'fullname': ['name', 'full name'],
        'full name': ['fullname', 'name'],
    }

    for alias in alias_map.get(normalized_key, []):
        if alias in mapping:
            return True, mapping.get(alias, '')

    return False, ''


def replace_placeholders_in_text(text: str, mapping: Dict[str, object]) -> str:
    if not text:
        return text

    def _replacement(match):
        raw_key = match.group(1)
        found, value = resolve_placeholder_value(mapping, raw_key)
        if not found:
            return match.group(0)
        try:
            if pd.isna(value):
                value = ''
        except Exception:
            pass
        return str(value)

    return PLACEHOLDER_RE.sub(_replacement, text)


def replace_placeholders_in_paragraph(paragraph, mapping: Dict[str, object]) -> None:
    runs = list(paragraph.runs)
    if not runs:
        paragraph.text = replace_placeholders_in_text(paragraph.text, mapping)
        return

    original_text = ''.join(run.text or '' for run in runs)
    replaced_text = replace_placeholders_in_text(original_text, mapping)
    if replaced_text == original_text:
        return

    runs[0].text = replaced_text
    for run in runs[1:]:
        run.text = ''


def find_and_replace_in_shape(shape, mapping):
    # Replace placeholders in text frames
    if not shape.has_text_frame:
        return
    text_frame = shape.text_frame
    for p in text_frame.paragraphs:
        replace_placeholders_in_paragraph(p, mapping)


def replace_placeholders(prs: Presentation, mapping: dict):
    for slide in prs.slides:
        for shape in slide.shapes:
            # text shapes
            try:
                find_and_replace_in_shape(shape, mapping)
            except Exception:
                # some shapes may throw; ignore and continue
                continue
            # tables
            if shape.shape_type == 19:  # TABLE
                table = shape.table
                for r in range(len(table.rows)):
                    for c in range(len(table.columns)):
                        cell = table.cell(r, c)
                        for p in cell.text_frame.paragraphs:
                            replace_placeholders_in_paragraph(p, mapping)


def format_filename(template: str, row: pd.Series, date_now: datetime):
    # allow placeholders like {FirstName}, {Fullname}, {date}, {timestamp}
    # create a case-insensitive mapping of column headers
    mapping_ci = {str(k).strip().lower(): ('' if pd.isna(v) else v) for k, v in row.items()}
    mapping_ci['date'] = date_now.strftime('%m-%d-%Y')
    # timestamp: time component only to pair with {date}
    mapping_ci['timestamp'] = date_now.strftime('%H-%M-%S-%f')

    # default naming when no explicit template is provided
    if not template:
        # Prefer common fullname columns; fallback to first column value
        name_val = (
            mapping_ci.get('fullname')
            or mapping_ci.get('full name')
            or mapping_ci.get('name')
        )
        if not name_val:
            try:
                name_val = row.iloc[0]
            except Exception:
                name_val = ''
        if pd.isna(name_val):
            name_val = ''
        filename = f"{name_val}_{mapping_ci['date']}_{mapping_ci['timestamp']}"
    else:
        filename = template
        # replace {Token} placeholders case-insensitively using mapping_ci
        for token in re.findall(r"\{([^{}]+)\}", filename):
            val = mapping_ci.get(token.strip().lower(), '')
            filename = filename.replace(f"{{{token}}}", str(val))

    filename = sanitize_filename(str(filename))
    return filename


def extract_recipient_email(row: pd.Series) -> str:
    for key, value in row.items():
        if 'email' in str(key).strip().lower() and not pd.isna(value):
            return str(value).strip()
    return ''


def pick_primary_output(out_paths: Dict, out_format: str) -> str:
    preferred = out_paths.get(out_format)
    if isinstance(preferred, list):
        return str(preferred[0]) if preferred else ''
    if preferred:
        return str(preferred)

    if out_paths.get('pdf'):
        return str(out_paths['pdf'])
    if out_paths.get('pptx'):
        return str(out_paths['pptx'])
    return ''


def resolve_soffice_path(soffice_path: str) -> Optional[str]:
    candidate = (soffice_path or '').strip()
    if candidate:
        expanded = os.path.expandvars(os.path.expanduser(candidate))
        if os.name == 'nt' and expanded.lower().endswith('soffice.exe'):
            alt_com = expanded[:-4] + '.com'
            if os.path.isfile(alt_com):
                expanded = alt_com
        if os.path.isfile(expanded):
            return expanded

        resolved = shutil.which(expanded)
        if resolved:
            return resolved

    if os.name == 'nt':
        win_candidates = []
        for env_key in ('PROGRAMFILES', 'PROGRAMFILES(X86)', 'PROGRAMW6432'):
            root = os.environ.get(env_key)
            if root:
                win_candidates.append(os.path.join(root, 'LibreOffice', 'program', 'soffice.com'))
                win_candidates.append(os.path.join(root, 'LibreOffice', 'program', 'soffice.exe'))

        win_candidates.extend([
            r'C:\Program Files\LibreOffice\program\soffice.com',
            r'C:\Program Files\LibreOffice\program\soffice.exe',
            r'C:\Program Files (x86)\LibreOffice\program\soffice.com',
            r'C:\Program Files (x86)\LibreOffice\program\soffice.exe',
        ])

        for item in win_candidates:
            if os.path.isfile(item):
                return item

        resolved_com = shutil.which('soffice.com')
        if resolved_com:
            return resolved_com

        resolved_exe = shutil.which('soffice.exe')
        if resolved_exe:
            return resolved_exe
    else:
        resolved = shutil.which('soffice')
        if resolved:
            return resolved

    return None


def to_file_uri(path_value: str) -> str:
    return pathlib.Path(path_value).resolve().as_uri()


def to_lo_user_installation_uri(path_value: str) -> str:
    normalized = pathlib.Path(path_value).resolve().as_posix()
    return f'file:///{normalized}'


def pptx_to_pdf(soffice_path: str, pptx_path: str, out_dir: str) -> str:
    # Convert pptx to pdf using LibreOffice soffice CLI
    resolved_soffice = resolve_soffice_path(soffice_path)
    if not resolved_soffice:
        raise FileNotFoundError(f"soffice not found at '{soffice_path}'. Please install LibreOffice or provide full path.")

    normalized_pptx_path = os.path.normpath(os.path.abspath(pptx_path))
    normalized_out_dir = os.path.normpath(os.path.abspath(out_dir))

    profile_dir = mkdtemp(prefix='lo-profile-')
    profile_uri = to_lo_user_installation_uri(profile_dir)

    cmd = [
        resolved_soffice,
        '--headless',
        '--nologo',
        '--nodefault',
        '--nofirststartwizard',
        '--norestore',
        '--nolockcheck',
        f'-env:UserInstallation={profile_uri}',
        '--convert-to',
        'pdf:impress_pdf_Export',
        '--outdir',
        normalized_out_dir,
        normalized_pptx_path,
    ]

    process_env = os.environ.copy()
    process_env.pop('PYTHONHOME', None)
    process_env.pop('PYTHONPATH', None)
    process_env['TEMP'] = profile_dir
    process_env['TMP'] = profile_dir

    try:
        subprocess.run(cmd, check=True, env=process_env, timeout=120)
    except subprocess.TimeoutExpired as exc:
        raise RuntimeError(f'LibreOffice conversion timed out after 120 seconds for: {pptx_path}') from exc
    finally:
        shutil.rmtree(profile_dir, ignore_errors=True)

    base = os.path.splitext(os.path.basename(pptx_path))[0]
    pdf_path = os.path.abspath(os.path.join(out_dir, base + '.pdf'))
    if not os.path.exists(pdf_path):
        raise FileNotFoundError('PDF conversion failed; expected output not found: ' + pdf_path)
    return pdf_path


def pptx_to_images_with_soffice(soffice_path: str, pptx_path: str, out_dir: str) -> List[str]:
    resolved_soffice = resolve_soffice_path(soffice_path)
    if not resolved_soffice:
        raise FileNotFoundError(f"soffice not found at '{soffice_path}'. Please install LibreOffice or provide full path.")

    normalized_pptx_path = os.path.normpath(os.path.abspath(pptx_path))
    normalized_out_dir = os.path.normpath(os.path.abspath(out_dir))
    profile_dir = mkdtemp(prefix='lo-profile-')
    profile_uri = to_lo_user_installation_uri(profile_dir)

    cmd = [
        resolved_soffice,
        '--headless',
        '--nologo',
        '--nodefault',
        '--nofirststartwizard',
        '--norestore',
        '--nolockcheck',
        f'-env:UserInstallation={profile_uri}',
        '--convert-to',
        'png',
        '--outdir',
        normalized_out_dir,
        normalized_pptx_path,
    ]

    process_env = os.environ.copy()
    process_env.pop('PYTHONHOME', None)
    process_env.pop('PYTHONPATH', None)
    process_env['TEMP'] = profile_dir
    process_env['TMP'] = profile_dir

    try:
        subprocess.run(cmd, check=True, env=process_env, timeout=120)
    finally:
        shutil.rmtree(profile_dir, ignore_errors=True)

    base = os.path.splitext(os.path.basename(pptx_path))[0]
    generated = sorted(pathlib.Path(out_dir).glob(f'{base}*.png'))
    return [str(path) for path in generated]


def pdf_to_images(pdf_path: str, dpi: int = 200, poppler_path: str = None):
    # returns list of PIL Images (one per page)
    images = convert_from_path(pdf_path, dpi=dpi, poppler_path=poppler_path)
    return images


def generate_for_row(prs_template: Presentation, row: pd.Series, out_dir: str, out_format: str,
                     filename_template: str, slide_indices: Optional[List[int]] = None,
                     soffice_path: str = 'soffice', poppler_path: str = None,
                     dpi: int = 200) -> Tuple[bool, str, Optional[Dict]]:
    """Generate certificate for a row. Returns (success, message, output_paths dict)."""
    date_now = datetime.now()
    filename_base = format_filename(filename_template, row, date_now)
    os.makedirs(out_dir, exist_ok=True)

    # choose a unique base name to avoid overwriting when duplicates occur
    unique_base = ensure_unique_basename(out_dir, filename_base, 'pptx')

    # copy presentation
    prs = Presentation()
    # load from template path by serializing: easier to use deep copy via save/load
    # Instead, save prs_template to a temp file and reopen
    fd, tmp_pptx = mkstemp(suffix='.pptx')
    os.close(fd)
    try:
        prs_template.save(tmp_pptx)
        prs = Presentation(tmp_pptx)
    except Exception:
        prs = prs_template

    # If slide_indices provided, remove all other slides so output contains only relevant slide(s)
    if slide_indices is not None:
        total = len(prs.slides)
        # normalize and keep unique valid indices
        keep = sorted(set(i for i in slide_indices if 0 <= i < total))
        # remove slides not in keep, iterate in reverse to avoid index shifting
        for i in range(total - 1, -1, -1):
            if i not in keep:
                try:
                    remove_slide_at_index(prs, i)
                except Exception:
                    # ignore removal errors and continue
                    continue

    # mapping: case-insensitive column headers -> values
    mapping = {str(k).strip().lower(): v for k, v in row.items()}
    replace_placeholders(prs, mapping)

    out_pptx = os.path.join(out_dir, unique_base + '.pptx')
    prs.save(out_pptx)

    out_paths = {}
    out_paths['pptx'] = out_pptx

    if out_format == 'pptx':
        return True, f"Generated: {out_paths}", out_paths

    # convert to pdf
    try:
        try:
            out_pdf = os.path.join(out_dir, unique_base + '.pdf')
            pdf_path = pptx_to_pdf(soffice_path, out_pptx, out_dir)
            if os.path.abspath(pdf_path) != os.path.abspath(out_pdf):
                shutil.move(pdf_path, out_pdf)
            out_paths['pdf'] = out_pdf
        except FileNotFoundError as e:
            # likely soffice not found; warn and skip conversion, keep PPTX only
            msg = f"PPTX->PDF conversion skipped for '{unique_base}': {e}"
            return False, msg, out_paths
        except subprocess.CalledProcessError as e:
            details = (e.stderr or e.output or '').strip()
            msg = f"PPTX->PDF conversion failed for '{unique_base}': {details or str(e)}"
            return False, msg, out_paths
        except Exception as e:
            msg = f"Unexpected error during PPTX->PDF conversion for '{unique_base}': {e}"
            return False, msg, out_paths

        if out_format in ('png', 'jpg', 'jpeg'):
            try:
                generated_png_paths = pptx_to_images_with_soffice(soffice_path, out_pptx, out_dir)

                if not generated_png_paths:
                    raise FileNotFoundError('LibreOffice image conversion did not produce PNG output.')

                if out_format == 'png':
                    if len(generated_png_paths) == 1:
                        source_path = generated_png_paths[0]
                        target_path = os.path.join(out_dir, unique_base + '.png')
                        if os.path.abspath(source_path) != os.path.abspath(target_path):
                            shutil.move(source_path, target_path)
                        out_paths['png'] = target_path
                    else:
                        saved = []
                        for index, source_path in enumerate(generated_png_paths, start=1):
                            target_path = os.path.join(out_dir, f"{unique_base}_page{index}.png")
                            if os.path.abspath(source_path) != os.path.abspath(target_path):
                                shutil.move(source_path, target_path)
                            saved.append(target_path)
                        out_paths['png'] = saved
                else:
                    saved = []
                    for index, source_path in enumerate(generated_png_paths, start=1):
                        img = Image.open(source_path)
                        target_path = os.path.join(
                            out_dir,
                            unique_base + '.jpg' if len(generated_png_paths) == 1 else f"{unique_base}_page{index}.jpg"
                        )
                        img.convert('RGB').save(target_path, 'JPEG')
                        saved.append(target_path)
                        try:
                            os.remove(source_path)
                        except Exception:
                            pass

                    out_paths['jpg'] = saved[0] if len(saved) == 1 else saved
            except Exception as image_error:
                msg = (
                    f"PPTX->image conversion failed for '{unique_base}': {image_error}. "
                    "If using legacy PDF-to-image fallback, install Poppler and add it to PATH."
                )
                return False, msg, out_paths
    finally:
        try:
            os.remove(tmp_pptx)
        except Exception:
            pass

    return True, f"Generated: {out_paths}", out_paths


def main(argv=None):
    parser = argparse.ArgumentParser(description='Generate certificates from a PPTX template and Excel data')
    parser.add_argument('--template', '-t', required=True, help='Path to PPTX template file')
    parser.add_argument('--data', '-d', required=True, help='Path to Excel file (.xlsx)')
    parser.add_argument('--outdir', '-o', required=True, help='Output directory')
    parser.add_argument('--format', '-f', choices=['pptx', 'pdf', 'png', 'jpg'], default='png', help='Export format')
    parser.add_argument('--name-template', '-n', default='', help='Filename template, use {ColumnName}, {date}, and {timestamp}. Default is Fullname_date_timestamp')
    parser.add_argument('--map', '-m', default='', help='Optional mapping from sheet name to slide index: "SheetA:0,SheetB:2" (0-based)')
    parser.add_argument('--soffice-path', default='soffice', help='Path to LibreOffice "soffice" executable (used for pptx->pdf)')
    parser.add_argument('--poppler-path', default=None, help='Optional path to poppler bin folder for pdf2image on Windows')
    parser.add_argument('--dpi', type=int, default=200, help='DPI for image rendering from PDF')
    parser.add_argument('--manifest', default='', help='Optional path to output JSON manifest for generated rows')
    args = parser.parse_args(argv)

    if not os.path.exists(args.template):
        print('Template file not found:', args.template)
        sys.exit(2)
    if not os.path.exists(args.data):
        print('Data file not found:', args.data)
        sys.exit(2)

    # Create date-based subfolder
    date_str = datetime.now().strftime('%Y-%m-%d')
    base_outdir = args.outdir
    args.outdir = os.path.join(base_outdir, date_str)
    os.makedirs(args.outdir, exist_ok=True)

    # Set up logging for failed items
    log_file = os.path.join(args.outdir, 'failed_items.log')
    logging.basicConfig(
        filename=log_file,
        level=logging.ERROR,
        format='%(asctime)s - %(message)s',
        datefmt='%Y-%m-%d %H:%M:%S'
    )

    # read data: read all sheets into a dict of DataFrames (or CSV as a single synthetic sheet)
    try:
        data_extension = os.path.splitext(args.data)[1].lower()
        if data_extension == '.csv':
            sheets = {'Responses': pd.read_csv(args.data)}
        else:
            sheets = pd.read_excel(args.data, sheet_name=None)
    except Exception as e:
        print('Failed to read data file:', e)
        sys.exit(3)

    # load template once
    try:
        prs_template = Presentation(args.template)
    except Exception as e:
        print('Failed to read template pptx:', e)
        sys.exit(4)

    explicit_map = parse_mapping_arg(args.map)

    sheet_names = list(sheets.keys())
    total_files = 0
    failed_count = 0
    manifest_records = []
    print(f'Generating files to "{args.outdir}" in format {args.format}...')

    for sheet_idx, sheet_name in enumerate(sheet_names):
        df = sheets[sheet_name]
        # determine slide indices for this sheet
        slide_indices: List[int] = []
        if sheet_name in explicit_map:
            slide_indices = [explicit_map[sheet_name]]
        else:
            slide_indices = find_slide_indices_by_sheetname(prs_template, sheet_name)
            if not slide_indices:
                # fallback: if counts match, map by order
                if len(sheet_names) == len(prs_template.slides):
                    slide_indices = [sheet_idx]
                else:
                    print(f'Warning: could not find slide for sheet "{sheet_name}" and no explicit map provided; skipping this sheet')
                    continue

        print(f'Processing sheet "{sheet_name}" with {len(df)} rows using slide indices {slide_indices}')
        for idx, row in df.iterrows():
            recipient_email = extract_recipient_email(row)
            try:
                success, message, res = generate_for_row(prs_template, row, args.outdir, args.format, args.name_template,
                                       slide_indices=slide_indices, soffice_path=args.soffice_path,
                                       poppler_path=args.poppler_path, dpi=args.dpi)
                if success:
                    print('Generated:', res)
                    total_files += 1
                    attachment_path = pick_primary_output(res or {}, args.format)
                    manifest_records.append({
                        'sheet_name': sheet_name,
                        'row_index': int(idx),
                        'recipient_email': recipient_email,
                        'status': 'success',
                        'error_message': None,
                        'attachment_path': attachment_path,
                        'filename': os.path.basename(attachment_path) if attachment_path else None,
                    })
                else:
                    print(f'Warning: {message}')
                    logging.error(f"Sheet: {sheet_name}, Row: {idx}, Error: {message}")
                    failed_count += 1
                    fallback_attachment = pick_primary_output(res or {}, args.format)
                    manifest_records.append({
                        'sheet_name': sheet_name,
                        'row_index': int(idx),
                        'recipient_email': recipient_email,
                        'status': 'fail',
                        'error_message': str(message),
                        'attachment_path': fallback_attachment,
                        'filename': os.path.basename(fallback_attachment) if fallback_attachment else None,
                    })
            except Exception as e:
                error_msg = f'Failed for sheet {sheet_name} row {idx}: {e}'
                print(error_msg)
                logging.error(error_msg)
                failed_count += 1
                manifest_records.append({
                    'sheet_name': sheet_name,
                    'row_index': int(idx),
                    'recipient_email': recipient_email,
                    'status': 'fail',
                    'error_message': str(e),
                    'attachment_path': None,
                    'filename': None,
                })

    summary = f'\nGeneration complete. Generated: {total_files} files, Failed: {failed_count} items.'
    print(summary)
    if failed_count > 0:
        print(f'See "{log_file}" for details on failed items.')
    logging.error(f'Generation Summary: {total_files} files generated, {failed_count} failed.')

    if args.manifest:
        try:
            manifest_dir = os.path.dirname(args.manifest)
            if manifest_dir:
                os.makedirs(manifest_dir, exist_ok=True)
            with open(args.manifest, 'w', encoding='utf-8') as manifest_file:
                json.dump(manifest_records, manifest_file, indent=2)
            print(f'Manifest written: {args.manifest}')
        except Exception as e:
            print(f'Warning: failed to write manifest: {e}')


if __name__ == '__main__':
    main()
