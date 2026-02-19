import sys
from pptx import Presentation
import os

def print_pptx_text(path):
    if not os.path.exists(path):
        print('File not found:', path)
        return
    prs = Presentation(path)
    print('---', path)
    for si, slide in enumerate(prs.slides):
        print(f'Slide {si}:')
        for shape in slide.shapes:
            try:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        print('  ', p.text)
            except Exception:
                continue

if __name__ == '__main__':
    if len(sys.argv) < 2:
        print('Usage: python verify_generated.py <out_dir>')
        sys.exit(2)
    out_dir = sys.argv[1]
    for name in os.listdir(out_dir):
        if name.lower().endswith('.pptx'):
            print_pptx_text(os.path.join(out_dir, name))

