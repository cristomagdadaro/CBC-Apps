"""
Create sample PPTX template and Excel data for testing the certificate generator.

Generates:
- sample_template.pptx
- sample_data.xlsx

Run with: python create_sample_files.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
import pandas as pd


def create_pptx(path: str):
    prs = Presentation()
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(4)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = 'Certificate of Completion'
    p.font.size = Pt(40)
    p.font.bold = True

    p = tf.add_paragraph()
    p.text = '<<FirstName>> <<Surname>>'
    p.font.size = Pt(32)

    p = tf.add_paragraph()
    p.text = 'Has successfully completed the course.'
    p.font.size = Pt(20)

    prs.save(path)
    print('Wrote sample template to', path)


def create_excel(path: str):
    df = pd.DataFrame([
        {'FirstName': 'John', 'Surname': 'Doe', 'GivenName': 'John', 'LastName': 'Doe'},
        {'FirstName': 'Jane', 'Surname': 'Smith', 'GivenName': 'Jane', 'LastName': 'Smith'},
    ])
    df.to_excel(path, index=False)
    print('Wrote sample data to', path)


if __name__ == '__main__':
    create_pptx('sample_template.pptx')
    create_excel('sample_data.xlsx')

